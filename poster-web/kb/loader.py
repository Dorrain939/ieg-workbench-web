"""知识库加载器：常见资料格式 → 可检索文本。

- txt / md：直读
- pdf：pypdf 抽取每页文本
- docx：python-docx 抽取段落 + 表格
- pptx：python-pptx 抽取每页文本框
- csv / tsv / xlsx / xls：抽取表格为 Markdown 表格
- png / jpg / webp / gif / bmp / svg：保存原图，并生成图片元信息文本

接口：load(file_path: Path) -> str（统一返回 utf-8 纯文本，段落用空行分隔）
"""
import csv
import pathlib
import xml.etree.ElementTree as ET


TEXT_EXTS = {".txt", ".md", ".markdown", ".json", ".html", ".htm"}
TABLE_EXTS = {".csv", ".tsv", ".xlsx", ".xls"}
IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".webp", ".gif", ".bmp", ".svg"}
SUPPORTED_EXTS = TEXT_EXTS | {".pdf", ".docx", ".pptx"} | TABLE_EXTS | IMAGE_EXTS


class LoaderError(Exception):
    pass


def load(file_path: pathlib.Path) -> str:
    """根据扩展名分发到对应解析器。返回纯文本。"""
    ext = file_path.suffix.lower()
    if ext in TEXT_EXTS:
        return _load_text(file_path)
    if ext in (".csv", ".tsv"):
        return _load_delimited(file_path, delimiter="\t" if ext == ".tsv" else ",")
    if ext in (".xlsx", ".xls"):
        return _load_spreadsheet(file_path)
    if ext in IMAGE_EXTS:
        return _load_image_meta(file_path)
    if ext == ".pdf":
        return _load_pdf(file_path)
    if ext == ".docx":
        return _load_docx(file_path)
    if ext == ".pptx":
        return _load_pptx(file_path)
    raise LoaderError(f"不支持的文件格式：{ext}")


def _load_text(fp: pathlib.Path) -> str:
    # 兜底多种编码
    for enc in ("utf-8", "utf-8-sig", "gbk", "gb18030"):
        try:
            return fp.read_text(encoding=enc)
        except UnicodeDecodeError:
            continue
    # 最后兜底
    return fp.read_text(encoding="utf-8", errors="replace")


def _markdown_table(rows: list[list[str]], max_rows: int = 500) -> str:
    clean_rows = []
    for row in rows[:max_rows]:
        cells = [str(c or "").replace("\n", " ").strip() for c in row]
        if any(cells):
            clean_rows.append(cells)
    if not clean_rows:
        return ""
    max_cols = min(max(len(r) for r in clean_rows), 20)
    norm = [(r + [""] * (max_cols - len(r)))[:max_cols] for r in clean_rows]
    header = norm[0]
    lines = ["| " + " | ".join(header) + " |"]
    lines.append("| " + " | ".join(["---"] * max_cols) + " |")
    for row in norm[1:]:
        lines.append("| " + " | ".join(row) + " |")
    if len(rows) > max_rows:
        lines.append(f"\n[表格过长，仅解析前 {max_rows} 行；原文件已完整保存。]")
    return "\n".join(lines)


def _load_delimited(fp: pathlib.Path, delimiter: str) -> str:
    text = _load_text(fp)
    rows = list(csv.reader(text.splitlines(), delimiter=delimiter))
    table = _markdown_table(rows)
    return f"[表格文件] {fp.name}\n\n{table}" if table else f"[表格文件] {fp.name}\n\n未识别到有效表格内容。"


def _load_spreadsheet(fp: pathlib.Path) -> str:
    ext = fp.suffix.lower()
    if ext == ".xls":
        try:
            import pandas as pd
            sheets = pd.read_excel(str(fp), sheet_name=None, header=None)
        except Exception as e:
            raise LoaderError(f"XLS 解析需要 pandas + xlrd 支持：{e}")
        parts = []
        for name, frame in sheets.items():
            rows = frame.fillna("").astype(str).values.tolist()
            table = _markdown_table(rows)
            if table:
                parts.append(f"[工作表] {name}\n{table}")
        return "\n\n".join(parts)

    try:
        from openpyxl import load_workbook
    except ImportError as e:
        raise LoaderError(f"XLSX 解析需要 openpyxl：{e}")
    wb = load_workbook(str(fp), read_only=True, data_only=True)
    parts = []
    for ws in wb.worksheets:
        rows = []
        for row in ws.iter_rows(values_only=True):
            rows.append([str(c or "").strip() for c in row])
        table = _markdown_table(rows)
        if table:
            parts.append(f"[工作表] {ws.title}\n{table}")
    return "\n\n".join(parts)


def _load_image_meta(fp: pathlib.Path) -> str:
    ext = fp.suffix.lower()
    if ext == ".svg":
        try:
            root = ET.fromstring(fp.read_text(encoding="utf-8", errors="replace"))
            texts = [node.text.strip() for node in root.iter() if node.text and node.text.strip()]
            return "\n".join([
                f"[图片文件] {fp.name}",
                "类型：SVG 矢量图",
                "可读文字：" + (" / ".join(texts[:40]) if texts else "未发现 SVG 文本节点"),
                "说明：原始图片已保存，可作为视觉参考或海报素材使用。",
            ])
        except Exception:
            return f"[图片文件] {fp.name}\n类型：SVG 矢量图\n说明：原始图片已保存，可作为视觉参考或海报素材使用。"
    try:
        from PIL import Image
    except ImportError as e:
        raise LoaderError(f"图片解析需要 Pillow：{e}")
    with Image.open(fp) as img:
        w, h = img.size
        mode = img.mode
        fmt = img.format or ext.lstrip(".").upper()
    return "\n".join([
        f"[图片文件] {fp.name}",
        f"格式：{fmt}",
        f"尺寸：{w}x{h}",
        f"色彩模式：{mode}",
        "说明：原始图片已保存，可作为 logo、头像、底图、装饰图、模块素材或视觉参考使用。",
    ])


def _load_pdf(fp: pathlib.Path) -> str:
    try:
        from pypdf import PdfReader
    except ImportError as e:
        raise LoaderError(f"PDF 解析需要 pypdf：{e}")
    reader = PdfReader(str(fp))
    if len(reader.pages) > 200:
        raise LoaderError(f"PDF 页数过多（{len(reader.pages)} > 200），请拆分后上传")
    parts = []
    for i, page in enumerate(reader.pages):
        try:
            txt = page.extract_text() or ""
        except Exception:
            txt = ""
        if txt.strip():
            parts.append(f"[第 {i+1} 页]\n{txt.strip()}")
    return "\n\n".join(parts)


def _iter_docx_blocks(doc):
    """Yield paragraphs/tables in their real Word order."""
    from docx.table import Table
    from docx.text.paragraph import Paragraph

    body = doc.element.body
    for child in body.iterchildren():
        if child.tag.endswith('}p'):
            yield Paragraph(child, doc)
        elif child.tag.endswith('}tbl'):
            yield Table(child, doc)


def _markdown_table_from_docx(table) -> str:
    rows = []
    for row in table.rows:
        cells = [c.text.strip().replace("\n", " ") for c in row.cells]
        # 保留空单元格位置，避免列错位；整行空才跳过。
        if any(cells):
            rows.append(cells)
    if not rows:
        return ""
    max_cols = max(len(r) for r in rows)
    norm = [r + [""] * (max_cols - len(r)) for r in rows]
    # 删除所有行都为空的尾列，避免合并单元格导出多余空列。
    while norm and norm[0] and all((row[-1] or "").strip() == "" for row in norm):
        norm = [row[:-1] for row in norm]
    if not norm or not norm[0]:
        return ""
    max_cols = len(norm[0])
    lines = ["| " + " | ".join(norm[0]) + " |"]
    lines.append("| " + " | ".join(["---"] * max_cols) + " |")
    for row in norm[1:]:
        lines.append("| " + " | ".join(row) + " |")
    return "\n".join(lines)


def _load_docx(fp: pathlib.Path) -> str:
    try:
        from docx import Document
    except ImportError as e:
        raise LoaderError(f"DOCX 解析需要 python-docx：{e}")
    doc = Document(str(fp))
    parts = []
    for block in _iter_docx_blocks(doc):
        if hasattr(block, "rows"):
            table_md = _markdown_table_from_docx(block)
            if table_md:
                parts.append(table_md)
        else:
            t = block.text.strip()
            if t:
                parts.append(t)
    return "\n\n".join(parts)


def _load_pptx(fp: pathlib.Path) -> str:
    try:
        from pptx import Presentation
    except ImportError as e:
        raise LoaderError(f"PPTX 解析需要 python-pptx：{e}")
    prs = Presentation(str(fp))
    parts = []
    for i, slide in enumerate(prs.slides):
        slide_parts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    txt = "".join(r.text for r in para.runs).strip()
                    if txt:
                        slide_parts.append(txt)
            elif shape.shape_type == 19:  # 表格
                if hasattr(shape, "table"):
                    for row in shape.table.rows:
                        cells = [c.text.strip() for c in row.cells if c.text.strip()]
                        if cells:
                            slide_parts.append(" | ".join(cells))
        if slide_parts:
            parts.append(f"[第 {i+1} 页]\n" + "\n".join(slide_parts))
    return "\n\n".join(parts)
